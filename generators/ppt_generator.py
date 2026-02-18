"""
generators/ppt_generator.py — Master PPTX generator using python-pptx.
Orchestrates slide rendering across all layout types.
Now theme-aware via PresentationTheme with consultant-quality visuals.
"""

from __future__ import annotations

import io
from pathlib import Path
from typing import Any, Dict, List, Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE

from agents.slide_content_agent import SlideContent
from config import OUTPUT_DIR, get_settings
from engine.pipeline_logger import PipelineLogger
from generators.chart_annotator import ChartAnnotator
from generators.exec_summary_builder import ExecSummaryBuilder
from generators.table_generator import TableGenerator
from generators.themes import PresentationTheme, THEME_CORPORATE_BLUE
from models import ChartData, SlidePlan, TableData


SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# ── Layout Position Constants (inches) ────────────────────
HEADER_LEFT = 0.7
HEADER_TOP = 0.3
HEADER_WIDTH = 12.0
HEADER_HEIGHT = 0.7
HEADER_LINE_TOP = 1.1

CONTENT_LEFT = 1.0
CONTENT_TOP = 1.4
CONTENT_WIDTH = 11.5
CONTENT_HEIGHT = 4.6

SPLIT_LEFT_WIDTH = 5.2

CHART_FULL_LEFT = 1.5
CHART_FULL_TOP = 1.4
CHART_FULL_WIDTH = 10.0
CHART_FULL_HEIGHT = 4.8

CHART_SPLIT_LEFT = 6.8
CHART_SPLIT_TOP = 1.6
CHART_SPLIT_WIDTH = 5.8
CHART_SPLIT_HEIGHT = 4.2

TABLE_LEFT = 1.0
TABLE_TOP = 1.4
TABLE_WIDTH = 11.0
TABLE_HEIGHT = 4.6

INSIGHT_BAR_LEFT = 0.5
INSIGHT_BAR_TOP = 6.15
INSIGHT_BAR_WIDTH = 12.3
INSIGHT_BAR_HEIGHT = 0.65

FOOTER_LEFT = 0.7
FOOTER_TOP = 7.0
FOOTER_WIDTH = 12.0
FOOTER_HEIGHT = 0.3


class InteractivePPTGenerator:
    """Builds a complete PPTX presentation from slide plans and content."""

    def __init__(self, theme: Optional[PresentationTheme] = None) -> None:
        self._settings = get_settings()
        self._log = PipelineLogger("PPTGenerator")
        self._theme = theme or THEME_CORPORATE_BLUE
        self._chart_gen = ChartAnnotator(theme=self._theme)
        self._table_gen = TableGenerator(theme=self._theme)
        self._exec_builder = ExecSummaryBuilder(theme=self._theme)
        self._prs: Optional[Presentation] = None
        self._total_slides: int = 0
        self._topic: str = ""

    @property
    def theme(self) -> PresentationTheme:
        return self._theme

    @theme.setter
    def theme(self, t: PresentationTheme) -> None:
        self._theme = t
        self._chart_gen.theme = t
        self._table_gen.theme = t
        self._exec_builder.theme = t

    def create_presentation(
        self,
        topic: str,
        slides: List[SlidePlan],
        contents: List[SlideContent],
        output_filename: Optional[str] = None,
    ) -> Path:
        """Generate a complete PPTX file."""
        self._log.action("Create Presentation", f"topic={topic[:50]}, slides={len(slides)}")

        self._prs = Presentation()
        self._prs.slide_width = SLIDE_WIDTH
        self._prs.slide_height = SLIDE_HEIGHT
        self._total_slides = len(slides)
        self._topic = topic

        for slide_plan, content in zip(slides, contents):
            with self._log.step_start(f"Slide {slide_plan.id}: {slide_plan.title[:30]}"):
                self._render_slide(slide_plan, content)

        if not output_filename:
            safe_topic = "".join(c if c.isalnum() or c in " -_" else "" for c in topic)
            output_filename = f"{safe_topic[:50].strip()}.pptx"

        output_path = OUTPUT_DIR / output_filename
        self._prs.save(str(output_path))
        self._log.info(f"Presentation saved: {output_path}")
        return output_path

    # ── Slide Dispatcher ────────────────────────────────────

    def _render_slide(self, plan: SlidePlan, content: SlideContent) -> None:
        # Full-slide image takes priority (from Render Deciding Agent)
        if getattr(content, "full_slide_image", None):
            self._render_full_image_slide(plan, content)
            return

        layout_type = plan.layout_type.lower()

        dispatch = {
            "title": self._render_title_slide,
            "bullet": self._render_bullet_slide,
            "chart": self._render_chart_slide,
            "table": self._render_table_slide,
            "split": self._render_split_slide,
            "exec_summary": self._render_exec_summary_slide,
            "section_divider": self._render_section_divider_slide,
            "closing": self._render_closing_slide,
        }

        renderer = dispatch.get(layout_type, self._render_bullet_slide)
        renderer(plan, content)

    def _render_full_image_slide(self, plan: SlidePlan, content: SlideContent) -> None:
        """Render a slide where a pre-rendered image covers the entire slide edge-to-edge."""
        self._log.info(f"Slide {plan.id}: rendering as full-slide image (from Render Deciding Agent)")
        slide = self._add_blank_slide()
        image_stream = io.BytesIO(content.full_slide_image)
        slide.shapes.add_picture(
            image_stream,
            Inches(0), Inches(0),
            SLIDE_WIDTH, SLIDE_HEIGHT,
        )

    # ── Decorative Helpers ───────────────────────────────────

    def _apply_gradient_background(self, slide, color_start: RGBColor,
                                    color_end: RGBColor, angle: int = 270) -> None:
        """Apply a two-stop gradient fill to the slide background."""
        bg = slide.background
        fill = bg.fill
        fill.gradient()
        fill.gradient_angle = float(angle)  # degrees (0=L→R, 90=B→T, 270=T→B)
        stops = fill.gradient_stops
        stops[0].color.rgb = color_start
        stops[0].position = 0.0
        stops[1].color.rgb = color_end
        stops[1].position = 1.0

    def _add_left_accent_bar(self, slide) -> None:
        """Add a thin vertical accent bar on the left edge of a content slide."""
        t = self._theme
        bar = slide.shapes.add_shape(
            1,  # MSO_SHAPE.RECTANGLE
            Inches(0), Inches(0), Inches(0.15), SLIDE_HEIGHT,
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = t.accent
        bar.line.fill.background()

    def _add_bottom_accent_strip(self, slide) -> None:
        """Add a thin accent-colored strip at the very bottom of the slide."""
        t = self._theme
        strip = slide.shapes.add_shape(
            1,  # MSO_SHAPE.RECTANGLE
            Inches(0), Inches(7.38), SLIDE_WIDTH, Inches(0.12),
        )
        strip.fill.solid()
        strip.fill.fore_color.rgb = t.primary
        strip.line.fill.background()

    def _add_footer_accent_line(self, slide) -> None:
        """Add a thin accent line just above the footer area."""
        t = self._theme
        line = slide.shapes.add_shape(
            1,  # MSO_SHAPE.RECTANGLE
            Inches(0.7), Inches(6.9), Inches(11.9), Inches(0.015),
        )
        line.fill.solid()
        line.fill.fore_color.rgb = t.text_muted
        line.line.fill.background()

    # ── Slide Renderers ─────────────────────────────────────

    def _render_title_slide(self, plan: SlidePlan, content: SlideContent) -> None:
        t = self._theme
        slide = self._add_blank_slide()

        # Gradient background
        self._apply_gradient_background(slide, t.primary, t.gradient_end, angle=315)

        WHITE = RGBColor(0xFF, 0xFF, 0xFF)

        # Top accent bar
        top_bar = slide.shapes.add_shape(
            1, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06),
        )
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = t.accent
        top_bar.line.fill.background()

        # Title (large, centered)
        title_box = slide.shapes.add_textbox(
            Inches(1.0), Inches(1.8), Inches(11.3), Inches(2.0)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = content.title
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = t.heading_font
        p.alignment = PP_ALIGN.CENTER

        # Accent divider (rounded rectangle)
        divider = slide.shapes.add_shape(
            5,  # MSO_SHAPE.ROUNDED_RECTANGLE
            Inches(4.5), Inches(4.0), Inches(4.3), Inches(0.06),
        )
        divider.fill.solid()
        divider.fill.fore_color.rgb = t.accent
        divider.line.fill.background()

        # Subtitle
        if content.key_insight:
            sub_box = slide.shapes.add_textbox(
                Inches(2.0), Inches(4.3), Inches(9.3), Inches(1.2)
            )
            tf = sub_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = content.key_insight
            p.font.size = Pt(20)
            p.font.color.rgb = WHITE
            p.font.name = t.body_font
            p.alignment = PP_ALIGN.CENTER

        # Bottom accent strip
        bottom_strip = slide.shapes.add_shape(
            1, Inches(0), Inches(7.38), SLIDE_WIDTH, Inches(0.12),
        )
        bottom_strip.fill.solid()
        bottom_strip.fill.fore_color.rgb = t.accent
        bottom_strip.line.fill.background()

    def _render_bullet_slide(self, plan: SlidePlan, content: SlideContent) -> None:
        t = self._theme
        slide = self._add_blank_slide()
        self._add_left_accent_bar(slide)
        self._add_bottom_accent_strip(slide)
        self._add_slide_header(slide, content.title, plan.id)

        if getattr(content, "infographic_image", None):
            self._add_infographic_image(slide, content.infographic_image, "full-slide")
        else:
            txBox = slide.shapes.add_textbox(
                Inches(CONTENT_LEFT), Inches(CONTENT_TOP),
                Inches(CONTENT_WIDTH), Inches(CONTENT_HEIGHT),
            )
            tf = txBox.text_frame
            tf.word_wrap = True
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

            for i, bullet in enumerate(content.content_bullets):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()

                # Accent-colored bullet character
                bullet_run = p.add_run()
                bullet_run.text = "\u25cf  "
                bullet_run.font.size = Pt(10)
                bullet_run.font.color.rgb = t.accent

                # Body text
                text_run = p.add_run()
                text_run.text = bullet
                text_run.font.size = Pt(16)
                text_run.font.color.rgb = t.text_dark
                text_run.font.name = t.body_font

                p.space_before = Pt(12)
                p.space_after = Pt(6)
                p.level = 0

        if content.key_insight:
            self._add_insight_bar(slide, content.key_insight)

        self._add_footer_accent_line(slide)
        self._add_footer(slide, plan.id)

    def _render_chart_slide(self, plan: SlidePlan, content: SlideContent) -> None:
        slide = self._add_blank_slide()
        self._add_left_accent_bar(slide)
        self._add_bottom_accent_strip(slide)
        self._add_slide_header(slide, content.title, plan.id)

        if content.chart_data:
            chart_buf = self._chart_gen.generate(content.chart_data)
            slide.shapes.add_picture(
                chart_buf,
                Inches(CHART_FULL_LEFT), Inches(CHART_FULL_TOP),
                Inches(CHART_FULL_WIDTH), Inches(CHART_FULL_HEIGHT),
            )
        else:
            self._log.warning(f"No chart data for slide {plan.id}, falling back to bullets")
            self._render_bullet_content(slide, content)

        if content.key_insight and not content.chart_data:
            self._add_insight_bar(slide, content.key_insight)

        self._add_footer_accent_line(slide)
        self._add_footer(slide, plan.id)

    def _render_table_slide(self, plan: SlidePlan, content: SlideContent) -> None:
        slide = self._add_blank_slide()
        self._add_left_accent_bar(slide)
        self._add_bottom_accent_strip(slide)
        self._add_slide_header(slide, content.title, plan.id)

        if content.table_data:
            self._table_gen.add_table(
                slide, content.table_data,
                left=TABLE_LEFT, top=TABLE_TOP,
                width=TABLE_WIDTH, height=TABLE_HEIGHT,
            )
        else:
            self._log.warning(f"No table data for slide {plan.id}, falling back to bullets")
            self._render_bullet_content(slide, content)

        self._add_footer_accent_line(slide)
        self._add_footer(slide, plan.id)

    def _render_split_slide(self, plan: SlidePlan, content: SlideContent) -> None:
        t = self._theme
        slide = self._add_blank_slide()
        self._add_left_accent_bar(slide)
        self._add_bottom_accent_strip(slide)
        self._add_slide_header(slide, content.title, plan.id)

        # Left panel — bullets with triangle markers (constrained to left of divider)
        # Use wider text box starting closer to left edge, with auto-shrink to prevent clipping
        left_panel_left = 0.5   # closer to left edge (after the 0.15" accent bar)
        left_panel_width = 5.7  # generous width, stops before the divider at 6.5"
        txBox = slide.shapes.add_textbox(
            Inches(left_panel_left), Inches(CONTENT_TOP),
            Inches(left_panel_width), Inches(CONTENT_HEIGHT),
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

        # Dynamically pick font size based on total bullet text length
        total_text = sum(len(b) for b in content.content_bullets)
        num_bullets = len(content.content_bullets)
        if total_text > 400 or num_bullets > 5:
            body_font_size = Pt(11)
            bullet_font_size = Pt(8)
            space_before = Pt(6)
            space_after = Pt(3)
        elif total_text > 250 or num_bullets > 4:
            body_font_size = Pt(12)
            bullet_font_size = Pt(9)
            space_before = Pt(7)
            space_after = Pt(3)
        else:
            body_font_size = Pt(14)
            bullet_font_size = Pt(10)
            space_before = Pt(8)
            space_after = Pt(4)

        for i, bullet in enumerate(content.content_bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()

            bullet_run = p.add_run()
            bullet_run.text = "\u25b8  "
            bullet_run.font.size = bullet_font_size
            bullet_run.font.color.rgb = t.accent

            text_run = p.add_run()
            text_run.text = bullet
            text_run.font.size = body_font_size
            text_run.font.color.rgb = t.text_dark
            text_run.font.name = t.body_font

            p.space_before = space_before
            p.space_after = space_after
            p.level = 0

        # Vertical divider line
        divider = slide.shapes.add_shape(
            1,  # MSO_SHAPE.RECTANGLE
            Inches(6.5), Inches(1.5), Inches(0.02), Inches(4.5),
        )
        divider.fill.solid()
        divider.fill.fore_color.rgb = t.accent
        divider.line.fill.background()

        # Right panel — chart, table, or image
        if content.chart_data:
            chart_buf = self._chart_gen.generate(content.chart_data)
            slide.shapes.add_picture(
                chart_buf,
                Inches(CHART_SPLIT_LEFT), Inches(CHART_SPLIT_TOP),
                Inches(CHART_SPLIT_WIDTH), Inches(CHART_SPLIT_HEIGHT),
            )
        elif content.table_data:
            self._table_gen.add_table(
                slide, content.table_data,
                left=CHART_SPLIT_LEFT, top=TABLE_TOP,
                width=CHART_SPLIT_WIDTH, height=CHART_SPLIT_HEIGHT,
            )
        elif getattr(content, "infographic_image", None):
            self._add_infographic_image(slide, content.infographic_image, "right-column")

        if content.key_insight:
            self._add_insight_bar(slide, content.key_insight)

        self._add_footer_accent_line(slide)
        self._add_footer(slide, plan.id)

    def _render_exec_summary_slide(self, plan: SlidePlan, content: SlideContent) -> None:
        slide = self._add_blank_slide()
        self._add_left_accent_bar(slide)
        self._add_bottom_accent_strip(slide)

        kpi_cards = self._extract_kpis(content.content_bullets)
        key_findings = [b for b in content.content_bullets if ":" not in b or len(b) > 60]

        self._exec_builder.build(
            slide=slide,
            title=content.title,
            kpi_cards=kpi_cards,
            key_findings=key_findings,
            bottom_insight=content.key_insight,
        )

        self._add_footer_accent_line(slide)
        self._add_footer(slide, plan.id)

    def _render_section_divider_slide(self, plan: SlidePlan, content: SlideContent) -> None:
        """Section divider slide with gradient background and large title."""
        t = self._theme
        slide = self._add_blank_slide()

        # Gradient background (diagonal)
        self._apply_gradient_background(slide, t.primary, t.gradient_end, angle=315)

        WHITE = RGBColor(0xFF, 0xFF, 0xFF)

        # Left accent bar (shorter, decorative)
        bar = slide.shapes.add_shape(
            1,  # MSO_SHAPE.RECTANGLE
            Inches(0.8), Inches(2.0), Inches(0.06), Inches(3.5),
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = t.accent
        bar.line.fill.background()

        # Section number label
        section_num = plan.id
        num_box = slide.shapes.add_textbox(
            Inches(1.2), Inches(2.0), Inches(10.0), Inches(0.5)
        )
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"SECTION {section_num}"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = t.accent
        p.font.name = t.heading_font

        # Section title (large, left-aligned)
        title_box = slide.shapes.add_textbox(
            Inches(1.2), Inches(2.7), Inches(10.0), Inches(1.8)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = content.title
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = t.heading_font

        # Optional subtitle from key_insight
        if content.key_insight:
            sub_box = slide.shapes.add_textbox(
                Inches(1.2), Inches(4.6), Inches(10.0), Inches(1.0)
            )
            tf = sub_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = content.key_insight
            p.font.size = Pt(16)
            p.font.color.rgb = WHITE
            p.font.name = t.body_font

        # Bottom accent strip
        bottom_strip = slide.shapes.add_shape(
            1, Inches(0), Inches(7.38), SLIDE_WIDTH, Inches(0.12),
        )
        bottom_strip.fill.solid()
        bottom_strip.fill.fore_color.rgb = t.accent
        bottom_strip.line.fill.background()

    def _render_closing_slide(self, plan: SlidePlan, content: SlideContent) -> None:
        """Closing/thank-you slide with gradient background."""
        t = self._theme
        slide = self._add_blank_slide()

        # Gradient background
        self._apply_gradient_background(slide, t.primary, t.gradient_end, angle=270)

        WHITE = RGBColor(0xFF, 0xFF, 0xFF)

        # Top accent bar
        top_bar = slide.shapes.add_shape(
            1, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06),
        )
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = t.accent
        top_bar.line.fill.background()

        # Title — "Thank You" or custom title
        title_text = content.title if content.title else "Thank You"
        title_box = slide.shapes.add_textbox(
            Inches(1.0), Inches(2.2), Inches(11.3), Inches(1.5)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = t.heading_font
        p.alignment = PP_ALIGN.CENTER

        # Accent divider
        divider = slide.shapes.add_shape(
            5,  # MSO_SHAPE.ROUNDED_RECTANGLE
            Inches(5.0), Inches(3.9), Inches(3.3), Inches(0.06),
        )
        divider.fill.solid()
        divider.fill.fore_color.rgb = t.accent
        divider.line.fill.background()

        # Subtitle / contact info from key_insight
        if content.key_insight:
            sub_box = slide.shapes.add_textbox(
                Inches(2.0), Inches(4.2), Inches(9.3), Inches(1.5)
            )
            tf = sub_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = content.key_insight
            p.font.size = Pt(18)
            p.font.color.rgb = WHITE
            p.font.name = t.body_font
            p.alignment = PP_ALIGN.CENTER

        # Bottom accent strip
        bottom_strip = slide.shapes.add_shape(
            1, Inches(0), Inches(7.38), SLIDE_WIDTH, Inches(0.12),
        )
        bottom_strip.fill.solid()
        bottom_strip.fill.fore_color.rgb = t.accent
        bottom_strip.line.fill.background()

    # ── Shared Components ───────────────────────────────────

    def _add_blank_slide(self):
        layout = self._prs.slide_layouts[6]
        return self._prs.slides.add_slide(layout)

    def _add_slide_header(self, slide, title: str, slide_id: int) -> None:
        t = self._theme
        txBox = slide.shapes.add_textbox(
            Inches(HEADER_LEFT), Inches(HEADER_TOP),
            Inches(HEADER_WIDTH), Inches(HEADER_HEIGHT),
        )
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = t.primary
        p.font.name = t.heading_font

        # Accent underline (wider)
        line = slide.shapes.add_shape(
            1, Inches(HEADER_LEFT), Inches(HEADER_LINE_TOP),
            Inches(3.0), Inches(0.04),
        )
        line.fill.solid()
        line.fill.fore_color.rgb = t.accent
        line.line.fill.background()

    def _add_footer(self, slide, slide_id: int) -> None:
        t = self._theme

        # Topic name on the left
        if self._topic:
            topic_box = slide.shapes.add_textbox(
                Inches(FOOTER_LEFT), Inches(FOOTER_TOP),
                Inches(6.0), Inches(FOOTER_HEIGHT),
            )
            tf = topic_box.text_frame
            p = tf.paragraphs[0]
            display_topic = self._topic[:50] + "..." if len(self._topic) > 50 else self._topic
            p.text = display_topic
            p.font.size = Pt(8)
            p.font.color.rgb = t.text_muted
            p.font.name = t.body_font
            p.alignment = PP_ALIGN.LEFT

        # "Slide X of Y" on the right
        slide_box = slide.shapes.add_textbox(
            Inches(10.5), Inches(FOOTER_TOP),
            Inches(2.2), Inches(FOOTER_HEIGHT),
        )
        tf = slide_box.text_frame
        p = tf.paragraphs[0]
        total = self._total_slides if self._total_slides > 0 else slide_id
        p.text = f"Slide {slide_id} of {total}"
        p.font.size = Pt(8)
        p.font.color.rgb = t.text_muted
        p.font.name = t.body_font
        p.alignment = PP_ALIGN.RIGHT

    def _add_insight_bar(self, slide, insight: str) -> None:
        t = self._theme

        # Rounded rectangle background with accent border
        shape = slide.shapes.add_shape(
            5,  # MSO_SHAPE.ROUNDED_RECTANGLE
            Inches(INSIGHT_BAR_LEFT), Inches(INSIGHT_BAR_TOP),
            Inches(INSIGHT_BAR_WIDTH), Inches(INSIGHT_BAR_HEIGHT),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = t.insight_bg
        shape.line.color.rgb = t.accent
        shape.line.width = Pt(1)

        txBox = slide.shapes.add_textbox(
            Inches(CONTENT_LEFT), Inches(INSIGHT_BAR_TOP + 0.05),
            Inches(11.5), Inches(0.55),
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]

        # Star prefix
        star_run = p.add_run()
        star_run.text = "\u2605 "
        star_run.font.size = Pt(11)
        star_run.font.color.rgb = t.accent

        # Label
        label_run = p.add_run()
        label_run.text = "KEY INSIGHT: "
        label_run.font.size = Pt(11)
        label_run.font.bold = True
        label_run.font.color.rgb = t.accent
        label_run.font.name = t.body_font

        # Insight text
        text_run = p.add_run()
        text_run.text = insight
        text_run.font.size = Pt(11)
        text_run.font.color.rgb = t.text_dark
        text_run.font.name = t.body_font

    def _add_infographic_image(
        self,
        slide,
        image_bytes: bytes,
        placement: str = "full-slide",
    ) -> None:
        """Add an infographic image to the slide."""
        image_stream = io.BytesIO(image_bytes)

        if placement == "right-column":
            slide.shapes.add_picture(
                image_stream,
                Inches(CHART_SPLIT_LEFT), Inches(CHART_SPLIT_TOP),
                Inches(CHART_SPLIT_WIDTH), Inches(CHART_SPLIT_HEIGHT),
            )
        elif placement == "bottom-section":
            bottom_top = CONTENT_TOP + CONTENT_HEIGHT / 2
            slide.shapes.add_picture(
                image_stream,
                Inches(CONTENT_LEFT), Inches(bottom_top),
                Inches(CONTENT_WIDTH), Inches(CONTENT_HEIGHT / 2),
            )
        else:
            # full-slide (default)
            slide.shapes.add_picture(
                image_stream,
                Inches(CONTENT_LEFT), Inches(CONTENT_TOP),
                Inches(CONTENT_WIDTH), Inches(CONTENT_HEIGHT),
            )

    def _render_bullet_content(self, slide, content: SlideContent) -> None:
        """Fallback bullet renderer for slides that expected chart/table data."""
        t = self._theme
        txBox = slide.shapes.add_textbox(
            Inches(CONTENT_LEFT), Inches(CONTENT_TOP),
            Inches(CONTENT_WIDTH), Inches(CONTENT_HEIGHT),
        )
        tf = txBox.text_frame
        tf.word_wrap = True

        for i, bullet in enumerate(content.content_bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()

            bullet_run = p.add_run()
            bullet_run.text = "\u25cf  "
            bullet_run.font.size = Pt(10)
            bullet_run.font.color.rgb = t.accent

            text_run = p.add_run()
            text_run.text = bullet
            text_run.font.size = Pt(16)
            text_run.font.color.rgb = t.text_dark
            text_run.font.name = t.body_font

            p.space_before = Pt(12)
            p.space_after = Pt(6)

    @staticmethod
    def _extract_kpis(bullets: List[str]) -> List[Dict[str, str]]:
        kpis: List[Dict[str, str]] = []
        for bullet in bullets:
            if ":" in bullet and len(bullet) < 60:
                parts = bullet.split(":", 1)
                label = parts[0].strip()
                rest = parts[1].strip()

                change = ""
                if "(" in rest and ")" in rest:
                    val_part = rest[:rest.index("(")].strip()
                    change = rest[rest.index("(") + 1:rest.index(")")].strip()
                else:
                    val_part = rest

                kpis.append({"label": label, "value": val_part, "change": change})

        return kpis[:4]
