"""
generators/exec_summary_builder.py — Complex grid layout for Executive Summary slides.
Creates card-based layouts with KPI metrics and key findings.
Now theme-aware via PresentationTheme.
"""

from __future__ import annotations

from typing import Any, Dict, List, Optional, TYPE_CHECKING

from pptx.slide import Slide
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from engine.pipeline_logger import PipelineLogger

if TYPE_CHECKING:
    from generators.themes import PresentationTheme


# ── Fallback Style Constants ───────────────────────────────────
CARD_BG = RGBColor(0xF8, 0xF9, 0xFA)
CARD_BORDER = RGBColor(0xDE, 0xE2, 0xE6)
ACCENT_COLOR = RGBColor(0x4A, 0x90, 0xD9)
NAVY = RGBColor(0x1B, 0x2A, 0x4A)
DARK_TEXT = RGBColor(0x2C, 0x3E, 0x50)
MUTED_TEXT = RGBColor(0x7F, 0x8C, 0x8D)


class ExecSummaryBuilder:
    """Builds executive summary slides with card-based grid layouts."""

    def __init__(self, theme: Optional[PresentationTheme] = None) -> None:
        self._log = PipelineLogger("ExecSummaryBuilder")
        self._theme = theme

    @property
    def theme(self) -> Optional[PresentationTheme]:
        return self._theme

    @theme.setter
    def theme(self, t: PresentationTheme) -> None:
        self._theme = t

    def _colors(self):
        t = self._theme
        if t:
            return {
                "card_bg": t.bg_light,
                "card_border": t.text_muted,
                "accent": t.accent,
                "primary": t.primary,
                "text_dark": t.text_dark,
                "text_muted": t.text_muted,
            }
        return {
            "card_bg": CARD_BG,
            "card_border": CARD_BORDER,
            "accent": ACCENT_COLOR,
            "primary": NAVY,
            "text_dark": DARK_TEXT,
            "text_muted": MUTED_TEXT,
        }

    def build(
        self,
        slide: Slide,
        title: str,
        kpi_cards: List[Dict[str, str]],
        key_findings: List[str],
        bottom_insight: str = "",
    ) -> None:
        """Build an executive summary slide with KPI cards and findings."""
        self._log.action("Build Exec Summary", f"title={title}")

        self._add_title(slide, title)
        self._add_kpi_grid(slide, kpi_cards)
        self._add_key_findings(slide, key_findings)

        if bottom_insight:
            self._add_bottom_insight(slide, bottom_insight)

        self._log.info(
            f"Exec summary built: {len(kpi_cards)} KPIs, {len(key_findings)} findings"
        )

    def _add_title(self, slide: Slide, title: str) -> None:
        c = self._colors()
        txBox = slide.shapes.add_textbox(
            Inches(0.7), Inches(0.3), Inches(12.3), Inches(0.6)
        )
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = c["primary"]

    def _add_kpi_grid(self, slide: Slide, cards: List[Dict[str, str]]) -> None:
        if not cards:
            return

        n_cards = min(len(cards), 4)
        card_width = 2.8
        card_height = 1.5
        gap = 0.3
        total_width = n_cards * card_width + (n_cards - 1) * gap
        start_left = (13.333 - total_width) / 2

        for i, card in enumerate(cards[:n_cards]):
            left = start_left + i * (card_width + gap)
            self._add_single_kpi_card(
                slide,
                left=left,
                top=1.2,
                width=card_width,
                height=card_height,
                label=card.get("label", ""),
                value=card.get("value", ""),
                change=card.get("change", ""),
            )

    def _add_single_kpi_card(
        self,
        slide: Slide,
        left: float,
        top: float,
        width: float,
        height: float,
        label: str,
        value: str,
        change: str,
    ) -> None:
        c = self._colors()

        # Rounded rectangle card background
        shape = slide.shapes.add_shape(
            5,  # MSO_SHAPE.ROUNDED_RECTANGLE
            Inches(left), Inches(top), Inches(width), Inches(height),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = c["card_bg"]
        shape.line.color.rgb = c["card_border"]
        shape.line.width = Pt(1)

        # Accent-coloured top border strip on card
        accent_strip = slide.shapes.add_shape(
            5,  # MSO_SHAPE.ROUNDED_RECTANGLE
            Inches(left + 0.1), Inches(top), Inches(width - 0.2), Inches(0.05),
        )
        accent_strip.fill.solid()
        accent_strip.fill.fore_color.rgb = c["accent"]
        accent_strip.line.fill.background()

        # Value (large)
        val_box = slide.shapes.add_textbox(
            Inches(left + 0.15), Inches(top + 0.2),
            Inches(width - 0.3), Inches(0.7),
        )
        tf = val_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = value
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = c["primary"]
        p.alignment = PP_ALIGN.CENTER

        # Label (below value)
        label_box = slide.shapes.add_textbox(
            Inches(left + 0.15), Inches(top + 0.85),
            Inches(width - 0.3), Inches(0.35),
        )
        tf = label_box.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(11)
        p.font.color.rgb = c["text_muted"]
        p.alignment = PP_ALIGN.CENTER

        # Change indicator
        if change:
            change_box = slide.shapes.add_textbox(
                Inches(left + 0.15), Inches(top + 1.15),
                Inches(width - 0.3), Inches(0.25),
            )
            tf = change_box.text_frame
            p = tf.paragraphs[0]
            p.text = change
            p.font.size = Pt(10)
            p.font.bold = True
            is_positive = change.strip().startswith("+") or change.strip().startswith("\u2191")
            p.font.color.rgb = (
                RGBColor(0x27, 0xAE, 0x60) if is_positive
                else RGBColor(0xE7, 0x4C, 0x3C)
            )
            p.alignment = PP_ALIGN.CENTER

    def _add_key_findings(self, slide: Slide, findings: List[str]) -> None:
        if not findings:
            return

        c = self._colors()
        txBox = slide.shapes.add_textbox(
            Inches(1.0), Inches(3.0), Inches(11.5), Inches(3.5)
        )
        tf = txBox.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = "Key Findings"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = c["primary"]
        p.space_after = Pt(8)

        for finding in findings[:6]:
            p = tf.add_paragraph()
            bullet_run = p.add_run()
            bullet_run.text = "\u25cf  "
            bullet_run.font.size = Pt(8)
            bullet_run.font.color.rgb = c["accent"]

            text_run = p.add_run()
            text_run.text = finding
            text_run.font.size = Pt(12)
            text_run.font.color.rgb = c["text_dark"]
            p.space_before = Pt(6)
            p.level = 0

    def _add_bottom_insight(self, slide: Slide, insight: str) -> None:
        c = self._colors()

        shape = slide.shapes.add_shape(
            5,  # MSO_SHAPE.ROUNDED_RECTANGLE
            Inches(0.5), Inches(6.15), Inches(12.3), Inches(0.7),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = c["accent"]
        shape.line.fill.background()

        txBox = slide.shapes.add_textbox(
            Inches(0.8), Inches(6.20), Inches(11.7), Inches(0.6)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = insight
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.alignment = PP_ALIGN.CENTER
