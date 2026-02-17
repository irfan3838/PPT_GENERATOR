"""
generators/table_generator.py — Data table formatting and rendering for PPTX slides.
Creates professionally styled tables using python-pptx.
Now theme-aware via PresentationTheme.
"""

from __future__ import annotations

from typing import Optional, TYPE_CHECKING

from pptx.slide import Slide
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from engine.pipeline_logger import PipelineLogger
from models import TableData

if TYPE_CHECKING:
    from generators.themes import PresentationTheme


# ── Fallback Style Constants ───────────────────────────────────
HEADER_BG = RGBColor(0x1B, 0x2A, 0x4A)       # Navy
HEADER_TEXT = RGBColor(0xFF, 0xFF, 0xFF)       # White
ROW_EVEN_BG = RGBColor(0xF8, 0xF9, 0xFA)      # Light gray
ROW_ODD_BG = RGBColor(0xFF, 0xFF, 0xFF)        # White
CELL_TEXT = RGBColor(0x2C, 0x3E, 0x50)         # Dark blue-gray
SOURCE_COLOR = RGBColor(0x95, 0xA5, 0xA6)      # Muted


class TableGenerator:
    """Generates professionally styled data tables on PPTX slides."""

    def __init__(self, theme: Optional[PresentationTheme] = None) -> None:
        self._log = PipelineLogger("TableGenerator")
        self._theme = theme

    @property
    def theme(self) -> Optional[PresentationTheme]:
        return self._theme

    @theme.setter
    def theme(self, t: PresentationTheme) -> None:
        self._theme = t

    def _colors(self):
        t = self._theme
        if t:
            return {
                "header_bg": t.primary,
                "header_text": RGBColor(0xFF, 0xFF, 0xFF),
                "row_even": t.bg_light,
                "row_odd": t.bg_white,
                "cell_text": t.text_dark,
                "source": t.text_muted,
            }
        return {
            "header_bg": HEADER_BG,
            "header_text": HEADER_TEXT,
            "row_even": ROW_EVEN_BG,
            "row_odd": ROW_ODD_BG,
            "cell_text": CELL_TEXT,
            "source": SOURCE_COLOR,
        }

    def add_table(
        self,
        slide: Slide,
        table_data: TableData,
        left: Optional[float] = None,
        top: Optional[float] = None,
        width: Optional[float] = None,
        height: Optional[float] = None,
    ) -> None:
        """Add a styled data table to a slide."""
        self._log.action("Add Table", f"'{table_data.title}' ({len(table_data.rows)} rows)")

        c = self._colors()
        n_rows = len(table_data.rows) + 1  # +1 for header
        n_cols = len(table_data.headers)

        left_emu = Inches(left if left is not None else 1.0)
        top_emu = Inches(top if top is not None else 2.0)
        width_emu = Inches(width if width is not None else 11.0)
        height_emu = Inches(height if height is not None else 4.0)

        table_shape = slide.shapes.add_table(
            n_rows, n_cols, left_emu, top_emu, width_emu, height_emu
        )
        table = table_shape.table

        col_width = int(width_emu / n_cols)
        for col_idx in range(n_cols):
            table.columns[col_idx].width = col_width

        # ── Header Row ──────────────────────────────────────
        for col_idx, header_text in enumerate(table_data.headers):
            cell = table.cell(0, col_idx)
            cell.text = header_text
            self._style_cell(
                cell,
                font_size=Pt(11),
                font_bold=True,
                font_color=c["header_text"],
                fill_color=c["header_bg"],
                alignment=PP_ALIGN.CENTER,
            )

        # ── Data Rows ───────────────────────────────────────
        for row_idx, row_data in enumerate(table_data.rows):
            bg_color = c["row_even"] if row_idx % 2 == 0 else c["row_odd"]
            for col_idx, cell_value in enumerate(row_data):
                if col_idx >= n_cols:
                    break
                cell = table.cell(row_idx + 1, col_idx)
                cell.text = str(cell_value)
                self._style_cell(
                    cell,
                    font_size=Pt(10),
                    font_bold=False,
                    font_color=c["cell_text"],
                    fill_color=bg_color,
                    alignment=PP_ALIGN.LEFT if col_idx == 0 else PP_ALIGN.CENTER,
                )

        # ── Source Annotation ───────────────────────────────
        if table_data.source_annotation:
            txBox = slide.shapes.add_textbox(
                left_emu,
                top_emu + height_emu + Inches(0.1),
                width_emu,
                Inches(0.3),
            )
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = table_data.source_annotation
            p.font.size = Pt(8)
            p.font.italic = True
            p.font.color.rgb = c["source"]

        self._log.info(f"Table added: {n_rows - 1} data rows x {n_cols} columns")

    @staticmethod
    def _style_cell(
        cell,
        font_size,
        font_bold: bool,
        font_color: RGBColor,
        fill_color: RGBColor,
        alignment,
    ) -> None:
        """Apply consistent styling to a table cell."""
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = font_size
            paragraph.font.bold = font_bold
            paragraph.font.color.rgb = font_color
            paragraph.alignment = alignment

        cell.fill.solid()
        cell.fill.fore_color.rgb = fill_color

        cell.margin_left = Inches(0.08)
        cell.margin_right = Inches(0.08)
        cell.margin_top = Inches(0.05)
        cell.margin_bottom = Inches(0.05)
